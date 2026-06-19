from flask import request, render_template, redirect, url_for, flash
from werkzeug.utils import secure_filename
from dataclasses import asdict
import google.generativeai as genai
from flask import Flask, json, render_template, request, redirect
from flask_sqlalchemy import SQLAlchemy, session
from sqlalchemy import or_, text
from werkzeug.utils import secure_filename
from pptx import Presentation
from flask import send_file
import fitz
from dotenv import load_dotenv
import pytesseract
from PIL import Image
import io
from datetime import date, datetime, timedelta
import json
import os
from groq import Groq
import aspose.slides as slides
import re
from dataclasses import asdict


load_dotenv()
from dataclasses import dataclass

@dataclass
class CriterionResult:

    category: str

    name: str

    required: str

    actual: str

    eligible: bool
#pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
groq_client = Groq(api_key=os.environ.get("GROQ_API_KEY"))


app = Flask(__name__, static_folder="static")
app.secret_key = "trident_secret_key_2024"
UPLOAD_FOLDER = "static/uploads"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

app.config["SQLALCHEMY_DATABASE_URI"] = "sqlite:///products.db"
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False
db = SQLAlchemy(app)

class FinancialYear(db.Model):
    __tablename__ = "financial_year"
    id = db.Column(db.Integer, primary_key=True)
    company_id = db.Column(db.Integer, db.ForeignKey("company_profile.id"), nullable=False)
    financial_year_end = db.Column(db.Integer, nullable=False)  # e.g. 2024 = FY 2023-24
    annual_turnover_crore = db.Column(db.Float, nullable=False)
    net_worth_crore = db.Column(db.Float)
    working_capital_crore = db.Column(db.Float)

    company = db.relationship("CompanyProfile", backref="financial_years")

class CompanyProfile(db.Model):

    __tablename__ = "company_profile"

    id = db.Column(

        db.Integer,

        primary_key=True

    )

    electrical_contractor_license = db.Column(

        db.String(100)

    )

    contractor_class = db.Column(

        db.String(50)

    )

    established_year = db.Column(

        db.Integer

    )

    entity_type = db.Column(
        
        db.String(50)
        
    )   # Proprietary / Company / Partnership / JV / Society / Trust / HUF / LLP
    
    pan_number = db.Column(
        
        db.String(20)
        
    )

class WorkExperience(db.Model):
    __tablename__ = "work_experience"
    id = db.Column(db.Integer, primary_key=True)
    company_id = db.Column(db.Integer, db.ForeignKey("company_profile.id"), nullable=False)
    project_name = db.Column(db.String(200))
    work_value_crore = db.Column(db.Float, nullable=False)
    completion_date = db.Column(db.Date, nullable=False)
    is_substantially_completed = db.Column(db.Boolean, default=False)
    work_type = db.Column(db.String(50))
    equipment_combo = db.Column(db.String(100))
    voltage_class_kv = db.Column(db.Float)
    certificate_issuer_type = db.Column(db.String(30))
    issuer_avg_turnover_3yr_crore = db.Column(db.Float)
    issuer_listed_nse = db.Column(db.Boolean, default=False)
    issuer_listed_bse = db.Column(db.Boolean, default=False)
    issuer_incorporation_date = db.Column(db.Date)
    has_work_order_copy = db.Column(db.Boolean, default=False)
    has_boq = db.Column(db.Boolean, default=False)
    has_ca_certified_payment_details = db.Column(db.Boolean, default=False)
    has_tds_certificates = db.Column(db.Boolean, default=False)
    has_final_bill_copy = db.Column(db.Boolean, default=False)

    company = db.relationship("CompanyProfile", backref="work_experiences")

class CompanyDocument(db.Model):
    __tablename__ = "company_document"
    id = db.Column(db.Integer, primary_key=True)
    company_id = db.Column(db.Integer, db.ForeignKey("company_profile.id"), nullable=False)
    document_type = db.Column(db.String(50), nullable=False)
    has_document = db.Column(db.Boolean, default=False)
    document_number = db.Column(db.String(100))
    issuing_authority = db.Column(db.String(150))
    valid_from = db.Column(db.Date)
    valid_until = db.Column(db.Date)
    file_path = db.Column(db.String(255))

    company = db.relationship("CompanyProfile", backref="documents")

class TenderComplianceRequirement(db.Model):
    __tablename__ = "tender_compliance_requirement"
    id = db.Column(db.Integer, primary_key=True)
    tender_history_id = db.Column(db.Integer, db.ForeignKey("tender_history.id"), nullable=False)
    section = db.Column(db.String(100))          # Commercial-Compliance / Technical-Compliances / Undertakings / etc.
    description = db.Column(db.Text)
    documents_uploading = db.Column(db.String(30))   # "Mandatory" / "Optional" / "Not Allowed"
    requirement_type = db.Column(db.String(50))   # normalized tag, matched against CompanyDocument.document_type; null for boilerplate items

    tender = db.relationship("TenderHistory", backref="compliance_requirements")

class Product(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    category = db.Column(db.String(100))
    description = db.Column(db.Text)
    features = db.Column(db.Text)
    applications = db.Column(db.Text)
    image_path = db.Column(db.String(255))

class TenderHistory(db.Model):

    id = db.Column(
        db.Integer,
        primary_key=True
    )

    filename = db.Column(
        db.String(255),
        nullable=False
    )

    tender_name = db.Column(
        db.String(300)
    )

    tender_value_crore = db.Column(
        db.Float
    )

    summary = db.Column(
        db.Text
    )

    eligibility_verdict = db.Column(
        db.String(50)
    )

    verdict_reason = db.Column(
        db.Text
    )

    recommended_action = db.Column(
        db.Text
    )

    eligibility_criteria = db.Column(
        db.Text
    )

    matched_products = db.Column(
        db.Text
    )

    created_at = db.Column(
        db.DateTime,
        default=datetime.utcnow
    )

    tender_number = db.Column(
        
        db.String(100)
        
    )
    
    tender_closing_date = db.Column(
        
        db.DateTime
        
    )
    
    local_content_percent = db.Column(
        
        db.Float
        
    )   # manually declared per bid, not derived from company data


class Service(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(200), nullable=False)
    short_description = db.Column(db.Text)
    detailed_content = db.Column(db.Text)  # AI-generated detailed content (JSON)
    image_url = db.Column(db.String(500))
    category = db.Column(db.String(100))
    icon = db.Column(db.String(10))
    link = db.Column(db.String(500))


class ManufacturingProduct(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(200), nullable=False)
    short_description = db.Column(db.Text)
    detailed_content = db.Column(db.Text)  # AI-generated detailed content (JSON)
    image_url = db.Column(db.String(500))
    category = db.Column(db.String(100))
    link = db.Column(db.String(500))

ENTITY_CERT_RULES = {
    "PARTNERSHIP_JV_HUF_LLP_CERTIFICATE": {"Partnership", "JV", "HUF", "LLP"},
    "SOLE_PROPRIETOR_UNDERTAKING": {"Proprietary"},
    "ANNEXURE_V_A": {"Partnership", "JV", "Society", "Trust", "HUF", "LLP"},
}

@app.route("/")
def home():

    recent_products = Product.query.order_by(Product.id.desc()).limit(4).all()

    return render_template("home.html", products=recent_products)

@app.route("/favicon.ico")
def favicon():
    return app.send_static_file("favicon.png")

from sqlalchemy import or_

@app.route("/test")
def test():
    return app.send_static_file("favicon.png")
@app.route("/products")
def products():

    search = request.args.get("search", "")

    print("==========")
    print("SEARCH =", search)

    if search:
        all_products = Product.query.filter(
            or_(
                Product.name.ilike(f"%{search}%"),
                Product.category.ilike(f"%{search}%"),
                Product.description.ilike(f"%{search}%"),
            )
        ).all()
    else:
        all_products = Product.query.all()

    print("RESULTS =", len(all_products))
    print("==========")
    print("SEARCH =", search)
    print("RESULTS =", len(all_products))
    return render_template("products.html", products=all_products, search=search)


@app.route("/add-product", methods=["GET", "POST"])
def add_product():

    if request.method == "POST":

        image = request.files["image"]
        filename = None

        if image and image.filename != "":
            filename = secure_filename(image.filename)

            image.save(os.path.join(app.config["UPLOAD_FOLDER"], filename))

        product = Product(
            name=request.form["name"],
            category=request.form["category"],
            description=request.form["description"],
            features=request.form["features"],
            applications=request.form["applications"],
            image_path=filename,
        )

        db.session.add(product)
        db.session.commit()

        return redirect("/products")

    return render_template("add_product.html")


@app.route("/product/<int:id>")
def product_details(id):

    product = Product.query.get_or_404(id)

    return render_template("product_details.html", product=product)


@app.route("/edit-product/<int:id>", methods=["GET", "POST"])
def edit_product(id):

    product = Product.query.get_or_404(id)

    if request.method == "POST":

        product.name = request.form["name"]
        product.category = request.form["category"]
        product.description = request.form["description"]
        product.features = request.form["features"]
        product.applications = request.form["applications"]

        image = request.files.get("image")

        if image and image.filename != "":
            filename = secure_filename(image.filename)

            image.save(os.path.join(app.config["UPLOAD_FOLDER"], filename))

            product.image_path = filename

        db.session.commit()

        return redirect("/products")

    return render_template("edit_product.html", product=product)


@app.route("/generate-ppt/<int:id>")
def generate_ppt(id):

    product = Product.query.get_or_404(id)

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = product.name

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = product.category

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Description"
    slide2.placeholders[1].text = product.description

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Features"
    slide3.placeholders[1].text = product.features

    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Applications"
    slide4.placeholders[1].text = product.applications

    filename = f"{product.name}.pptx"

    prs.save(filename)

    return send_file(filename, as_attachment=True)

def extract_pdf_text(pdf):

    text = ""

    doc = fitz.open(
        stream=pdf.read(),

        filetype="pdf"
    )

    for page in doc:

        text += page.get_text()

    return text


def extract_sections(full_text):

    headings = [

        "1. NIT HEADER",

        "2. SCHEDULE",

        "Eligibility Conditions",

        "Special Financial Criteria",

        "Special Technical Criteria",

        "Commercial-Compliance",

        "General Instructions",

        "Special Conditions",

        "Technical-Compliances",

        "Undertakings",

        "Custom"

    ]

    lower = full_text.lower()

    positions = []

    for heading in headings:

        idx = lower.find(

            heading.lower()

        )

        if idx != -1:

            positions.append(

                (idx, heading)

            )

    positions.sort()

    sections = {}

    for i, (start, heading) in enumerate(positions):
        if i < len(positions)-1:
            end = positions[i+1][0]
        else:
            end = len(full_text)
        sections[heading] = \
            full_text[start:end]
    return sections

def build_dashboard_data(sections):

    nit_data = build_nit_dictionary(sections)

    advertised = nit_data.get(

        "Advertised Value",

        "-"

    )


    advertised = re.sub(
        r"[^\d.]",
        "",
        str(advertised)
    )

    if advertised:
        advertised = f"{float(advertised):,.2f}"
    else:
        advertised = "-"


    earnest = nit_data.get(

        "Earnest Money (Rs.)",

        "-"

    )


    if earnest != "-":

        earnest = f"{float(earnest):,.2f}"


    completion = nit_data.get(

        "Period of Completion",

        "-"

    )


    return {

        "advertised_value":

        advertised,


        "earnest_money":

        earnest,


        "completion_period":

        completion

    }

def build_nit_dictionary(sections):

    nit = sections.get(
        "1. NIT HEADER",
        ""
    )

    fields = [

        "Name of Work",

        "Tender Closing Date Time",

        "Bidding Start Date",

        "Bidding type",

        "Tender Type",

        "Bidding System",

        "Pre-Bid Conference Required",

        "Advertised Value",

        "Earnest Money (Rs.)",

        "Validity of Offer ( Days)",

        "Period of Completion",

        "Are JV allowed to bid",

        "Are Consortium allowed to bid",

        "Contract Type"

    ]

    nit_data = {}

    for field in fields:
        pattern = rf"{re.escape(field)}\s*:?\s*(.+)"
        match = re.search(pattern,nit,re.IGNORECASE)
        if match:
            nit_data[field] = \
                match.group(1).strip()
     # ---------- SPECIAL FIXES ----------

    consortium = re.search(

        r'Are Consortium allowed\s+to bid\s+(Yes|No)',

        nit,

        re.IGNORECASE

    )

    if consortium:
        nit_data["Are Consortium allowed to bid"] = \
            consortium.group(1)


    closing = re.search(

        r'Tender Closing Date\s*Time\s*(\d{2}/\d{2}/\d{4}\s+\d{2}:\d{2})',

        nit,

        re.IGNORECASE

    )

    if closing:
        nit_data["Tender Closing Date Time"] = \
            closing.group(1)


    completion = re.search(

        r'Period of Completion\s+(\d+\s+\w+)',

        nit,

        re.IGNORECASE

    )

    if completion: 
        nit_data["Period of Completion"] = \
            completion.group(1)
    return nit_data

def answer_question(

question,

dashboard,

nit_data

):

    Q_MAP = {

        "When does bidding start?":

        lambda:

        nit_data.get(

            "Bidding Start Date",

            "Not found"

        ),


        "When is the tender closing date?":

        lambda:

        nit_data.get(

            "Tender Closing Date Time",

            "Not found"

        ),


        "Is JV allowed?":

        lambda:

        nit_data.get(

            "Are JV allowed to bid",

            "Not found"

        ),


        "Is consortium allowed?":

        lambda:

        nit_data.get(

            "Are Consortium allowed to bid",

            "Not found"

        ),


        "What is the contract type?":

        lambda:

        nit_data.get(

            "Contract Type",

            "Not found"

        )

    }


    if question in Q_MAP:

        return Q_MAP[question]()


    return "Question unavailable."

def evaluate_financial_eligibility(requirements, profile):
    results = []
    recent = (FinancialYear.query.filter_by(company_id=profile.id)
              .order_by(FinancialYear.financial_year_end.desc())
              .limit(3).all())
    avg_turnover_rupees = (
        sum(fy.annual_turnover_crore for fy in recent) / len(recent) * 1e7
        if recent else 0
    )
    for req in requirements:
        if req["type"] == "turnover":
            results.append(CriterionResult(
                category="Financial", name="Average Annual Turnover (last 3 FY)",
                required=f"Rs. {req['required']:,.0f}",
                actual=f"Rs. {avg_turnover_rupees:,.0f}",
                eligible=avg_turnover_rupees >= req["required"]
            ))
    return results

def evaluate_technical_eligibility(technical, tender_value_crore):
    cutoff = date.today() - timedelta(days=365 * technical.get("lookback_period_years", 7))
    defn = technical.get("similar_work_definition", {})
    allowed_types = set(defn.get("allowed_work_types", []))
    allowed_combos = set(defn.get("allowed_equipment_combos", []))
    min_kv = defn.get("min_voltage_kv", 0)
    cert_rules = technical.get("certificate_rules", {})

    works = WorkExperience.query.filter(
        WorkExperience.completion_date >= cutoff,
        WorkExperience.is_substantially_completed == True
    ).all()

    qualifying = []
    for w in works:
        if allowed_types and w.work_type not in allowed_types:
            continue
        if allowed_combos and w.equipment_combo not in allowed_combos:
            continue
        if w.voltage_class_kv < min_kv:
            continue
        qualifying.append({"work": w, "percent": (w.work_value_crore / tender_value_crore) * 100})

    for opt in sorted(technical.get("similar_work_options", []), key=lambda o: -o["min_count"]):
        matches = [q for q in qualifying if q["percent"] >= opt["min_percent_of_tender_value"]]
        if len(matches) >= opt["min_count"]:
            return CriterionResult(
                category="Technical", name="Similar Work Experience",
                required=f"{opt['min_count']} work(s) ≥{opt['min_percent_of_tender_value']}% of tender value",
                actual=f"{len(matches)} qualifying work(s) found",
                eligible=True
            )

    best = max(qualifying, key=lambda q: q["percent"], default=None)
    return CriterionResult(
        category="Technical", name="Similar Work Experience",
        required="See similar_work_options thresholds",
        actual=f"Best match {best['percent']:.0f}% of tender value" if best else "No qualifying work found",
        eligible=False
    )

def evaluate_all_eligibility(sections, profile, tender_value_crore, technical_json):
    results = []
    results += evaluate_financial_eligibility(extract_financial_eligibility(sections), profile)
    results.append(evaluate_technical_eligibility(technical_json, tender_value_crore))

    verdict = "ELIGIBLE" if all(r.eligible for r in results) else "NOT ELIGIBLE"
    failed = [r for r in results if not r.eligible]
    reason = "; ".join(f"{r.name}: required {r.required}, actual {r.actual}" for r in failed) or "All criteria met"

    return {
        "verdict": verdict,
        "reason": reason,
        "criteria_json": json.dumps([asdict(r) for r in results])
    }

def evaluate_compliance_eligibility(requirements, profile):
    results = []
    docs = {d.document_type: d for d in CompanyDocument.query.filter_by(company_id=profile.id)}

    for req in requirements:
        if req.requirement_type is None or req.documents_uploading != "Mandatory":
            continue  # boilerplate undertaking or optional item — not a pass/fail gate

        applicable_entities = ENTITY_CERT_RULES.get(req.requirement_type)
        if applicable_entities and profile.entity_type not in applicable_entities:
            continue  # doesn't apply to this entity type

        doc = docs.get(req.requirement_type)
        has_valid_doc = bool(doc and doc.has_document and (
            doc.valid_until is None or doc.valid_until >= date.today()
        ))
        results.append(CriterionResult(
            category="Compliance", name=req.requirement_type,
            required="Document on file", actual="Present" if has_valid_doc else "Missing",
            eligible=has_valid_doc
        ))
    return results

def extract_eligibility(sections):
    """
    Extracts financial eligibility requirements from the
    'Special Financial Criteria' section of the tender.
    """
    requirements = []
    financial = sections.get("Special Financial Criteria", "")

    # Convert newlines/tabs/multiple spaces into one space
    financial = re.sub(r"\s+", " ", financial)

    # Find all turnover requirement mentions, capturing the unit too
    turnover_matches = re.findall(
        r"annual contractual turnover\s*=\s*Rs\.?\s*([\d,]+(?:\.\d+)?)\s*(crore|lakh)?",
        financial,
        re.IGNORECASE
    )

    if turnover_matches:
        value_str, unit = turnover_matches[-1]  # take the last occurrence
        required_turnover = float(value_str.replace(",", ""))

        if unit.lower() == "crore":
            required_turnover *= 1e7
        elif unit.lower() == "lakh":
            required_turnover *= 1e5
        # if no unit is found, value is assumed to already be in rupees

        requirements.append({
            "type": "turnover",
            "name": "Average Annual Turnover",
            "required": required_turnover
        })

    return requirements

def parse_gemini(response):

    if hasattr(response, "text"):

        response = response.text

    response = response.replace(
        "```json",
        ""
    )

    response = response.replace(
        "```",
        ""
    )

    response = response.strip()

    try:

        return json.loads(response)

    except Exception:

        return {

        "lookback_period_years":7,

        "similar_work_options":[],

        "similar_work_definition":{

            "allowed_work_types":[],

            "allowed_equipment_combos":[],

            "min_voltage_kv":0

        },

        "certificate_rules":{},

        "mandatory_conditions":[]

        }
    
def summarize_technical_criteria(technical):

    model = genai.GenerativeModel(
        "gemini-2.5-flash"
    )

    prompt = f"""
You are a tender eligibility extractor.

Analyze this tender technical criteria.

Extract:

- years_of_experience
- similar_work_options
- similar_work_definition
- mandatory_conditions

Return ONLY JSON.

Example:

{{
 "lookback_period_years":7,

 "similar_work_options":[

   {{
      "min_count":3,

      "min_percent_of_tender_value":30
   }},

   {{
      "min_count":2,

      "min_percent_of_tender_value":40
   }},

   {{
      "min_count":1,

      "min_percent_of_tender_value":60
   }}

 ],

 "similar_work_definition":{{

   "allowed_work_types":[],

   "allowed_equipment_combos":[],

   "min_voltage_kv":0

 }},

 "certificate_rules":{{}},

 "mandatory_conditions":[]
}}

Tender text:

{technical}

Return JSON only.
"""

    response = model.generate_content(
        prompt
    )

    return response.text

@app.route("/tender", methods=["GET", "POST"])
def tender():
    if request.method == "GET":
        return render_template("tender_result.html", result= None , history = None , criteria=[])

    pdf = request.files.get("pdf_file")
    if not pdf or pdf.filename == "":
        flash("Please upload a tender PDF.")
        return redirect(url_for("tender"))

    filename = secure_filename(pdf.filename)

    # extract_pdf_text reads via pdf.read(), so do that before saving
    full_text = extract_pdf_text(pdf)
    sections = extract_sections(full_text)
    nit_data = build_nit_dictionary(sections)

    # tender value, needed for technical % thresholds
    advertised_value_str = nit_data.get("Advertised Value", "0")
    try:
        tender_value_crore = float(re.sub(r"[^\d.]", "", advertised_value_str)) / 1e7
    except ValueError:
        tender_value_crore = 0.0

    # financial check

    financial_requirements = extract_eligibility(sections)

    profile = CompanyProfile.query.first()

    financial_results = evaluate_financial_eligibility(financial_requirements,profile)

    dashboard = build_dashboard_data(sections)


    # technical check

    technical_section_text = sections.get("Special Technical Criteria","")

    gemini_response = summarize_technical_criteria(technical_section_text)

    technical_json = parse_gemini(gemini_response)

    technical_result = evaluate_technical_eligibility(technical_json,tender_value_crore)



    # combine results

    results = financial_results + [technical_result]
    # compliance check (entity type / documents) — only runs if requirements were parsed
    # for this tender; see note below

    overall_eligible = all(r.eligible for r in results)
    failed = [r for r in results if not r.eligible]
    reason = "; ".join(f"{r.name}: required {r.required}, actual {r.actual}" for r in failed) or "All criteria met"

    history = TenderHistory(
        filename=filename,
        tender_name=nit_data.get("Name of Work", "-"),
        tender_value_crore=tender_value_crore,
        eligibility_verdict="ELIGIBLE" if overall_eligible else "NOT ELIGIBLE",
        verdict_reason=reason,
        eligibility_criteria=json.dumps([asdict(r) for r in results]),
        created_at=datetime.utcnow()
    )
    db.session.add(history)
    db.session.commit()
    result = {"nit_header": nit_data,"dashboard": dashboard, "eligibility": financial_results,"technical_eligibility": technical_json}
    return render_template(

        "tender_result.html",

        result=result,

        history=history,

        criteria=results

    )
@app.route("/services")
def services():
    all_services = Service.query.all()
    # Parse detailed_content JSON for each service
    for s in all_services:
        if s.detailed_content:
            try:
                s._parsed_content = json.loads(s.detailed_content)
            except:
                s._parsed_content = None
        else:
            s._parsed_content = None
    return render_template("services.html", services=all_services)


@app.route("/manufacturing")
def manufacturing():
    all_products = ManufacturingProduct.query.all()
    for p in all_products:
        if p.detailed_content:
            try:
                p._parsed_content = json.loads(p.detailed_content)
            except:
                p._parsed_content = None
        else:
            p._parsed_content = None
    return render_template("manufacturing.html", products=all_products)


@app.route("/service/<int:id>")
def service_detail(id):
    service = Service.query.get_or_404(id)
    parsed_content = None
    if service.detailed_content:
        try:
            parsed_content = json.loads(service.detailed_content)
        except:
            parsed_content = None
    return render_template("service_detail.html", service=service, content=parsed_content)


@app.route("/manufacturing-product/<int:id>")
def manufacturing_detail(id):
    product = ManufacturingProduct.query.get_or_404(id)
    parsed_content = None
    if product.detailed_content:
        try:
            parsed_content = json.loads(product.detailed_content)
        except:
            parsed_content = None
    return render_template("manufacturing_detail.html", product=product, content=parsed_content)


@app.route("/generate-service-content/<int:id>", methods=["POST"])
def generate_service_content(id):
    service = Service.query.get_or_404(id)

    prompt = f"""
You are a professional content writer for Trident Engineers & Associates, an electrical and instrumentation EPC company.

Generate detailed, professional content for this service:
Service Name: {service.name}
Short Description: {service.short_description}
Category: {service.category}

Create comprehensive content in this exact JSON format:
{{
    "tagline": "A compelling one-line tagline for this service",
    "overview": "A detailed 3-4 paragraph overview of this service, explaining what it involves, why it's important, and how Trident delivers it professionally. Be specific to electrical/industrial engineering.",
    "key_features": [
        "Feature 1 with brief explanation",
        "Feature 2 with brief explanation",
        "Feature 3 with brief explanation",
        "Feature 4 with brief explanation",
        "Feature 5 with brief explanation",
        "Feature 6 with brief explanation"
    ],
    "applications": [
        "Application/Industry 1",
        "Application/Industry 2",
        "Application/Industry 3",
        "Application/Industry 4",
        "Application/Industry 5"
    ],
    "process_steps": [
        {{"step": "Step 1 title", "description": "Brief description"}},
        {{"step": "Step 2 title", "description": "Brief description"}},
        {{"step": "Step 3 title", "description": "Brief description"}},
        {{"step": "Step 4 title", "description": "Brief description"}}
    ],
    "why_choose_us": [
        "Reason 1",
        "Reason 2",
        "Reason 3",
        "Reason 4"
    ],
    "technical_specs": "A paragraph about technical standards, certifications, and compliance relevant to this service"
}}

Make the content professional, technically accurate, and specific to the Indian industrial/electrical sector. Reference Indian standards (IS, BIS, CEA) where appropriate.
"""

    response = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
    )
    text_resp = (
        response.choices[0]
        .message.content.strip()
        .replace("```json", "")
        .replace("```", "")
    )

    try:
        content = json.loads(text_resp)
        service.detailed_content = json.dumps(content)
        db.session.commit()
    except:
        pass

    return redirect(f"/service/{id}")


@app.route("/generate-manufacturing-content/<int:id>", methods=["POST"])
def generate_manufacturing_content(id):
    product = ManufacturingProduct.query.get_or_404(id)

    prompt = f"""
You are a professional content writer for Trident Engineers & Associates, an electrical panel and industrial manufacturing company.

Generate detailed, professional content for this manufactured product:
Product Name: {product.name}
Short Description: {product.short_description}
Category: {product.category}

Create comprehensive content in this exact JSON format:
{{
    "tagline": "A compelling one-line tagline for this product",
    "overview": "A detailed 3-4 paragraph overview of this product, explaining what it is, its construction, materials used, and how Trident manufactures it. Be specific to electrical/industrial manufacturing.",
    "key_features": [
        "Feature 1 with brief explanation",
        "Feature 2 with brief explanation",
        "Feature 3 with brief explanation",
        "Feature 4 with brief explanation",
        "Feature 5 with brief explanation",
        "Feature 6 with brief explanation"
    ],
    "specifications": [
        {{"spec": "Specification name", "value": "Specification value/range"}},
        {{"spec": "Specification name", "value": "Specification value/range"}},
        {{"spec": "Specification name", "value": "Specification value/range"}},
        {{"spec": "Specification name", "value": "Specification value/range"}},
        {{"spec": "Specification name", "value": "Specification value/range"}}
    ],
    "applications": [
        "Application/Industry 1",
        "Application/Industry 2",
        "Application/Industry 3",
        "Application/Industry 4",
        "Application/Industry 5"
    ],
    "manufacturing_process": [
        {{"step": "Step 1 title", "description": "Brief description"}},
        {{"step": "Step 2 title", "description": "Brief description"}},
        {{"step": "Step 3 title", "description": "Brief description"}},
        {{"step": "Step 4 title", "description": "Brief description"}}
    ],
    "quality_standards": "A paragraph about quality standards, testing procedures, certifications, and compliance relevant to this product. Reference Indian standards (IS, BIS, IEC) where appropriate.",
    "why_choose_us": [
        "Reason 1",
        "Reason 2",
        "Reason 3",
        "Reason 4"
    ]
}}

Make the content professional, technically accurate, and specific to the Indian industrial/electrical manufacturing sector.
"""

    response = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
    )
    text_resp = (
        response.choices[0]
        .message.content.strip()
        .replace("```json", "")
        .replace("```", "")
    )

    try:
        content = json.loads(text_resp)
        product.detailed_content = json.dumps(content)
        db.session.commit()
    except:
        pass

    return redirect(f"/manufacturing-product/{id}")


@app.route("/generate-all-content", methods=["POST"])
def generate_all_content():
    """Generate content for all services and manufacturing products that don't have content yet."""
    item_type = request.form.get("type", "all")

    generated = 0

    if item_type in ["all", "services"]:
        services = Service.query.filter(
            (Service.detailed_content == None) | (Service.detailed_content == "")
        ).all()
        for service in services:
            try:
                prompt = f"""You are a professional content writer for Trident Engineers & Associates, an electrical and instrumentation EPC company.
Generate detailed content for: {service.name}
Description: {service.short_description}
Category: {service.category}

Respond ONLY in this exact JSON format:
{{"tagline": "compelling tagline", "overview": "3-4 detailed paragraphs about this service", "key_features": ["feature1", "feature2", "feature3", "feature4", "feature5", "feature6"], "applications": ["app1", "app2", "app3", "app4", "app5"], "process_steps": [{{"step": "title", "description": "desc"}}, {{"step": "title", "description": "desc"}}, {{"step": "title", "description": "desc"}}, {{"step": "title", "description": "desc"}}], "why_choose_us": ["reason1", "reason2", "reason3", "reason4"], "technical_specs": "paragraph about standards and compliance"}}"""

                response = groq_client.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role": "user", "content": prompt}],
                )
                text_resp = (
                    response.choices[0]
                    .message.content.strip()
                    .replace("```json", "")
                    .replace("```", "")
                )
                content = json.loads(text_resp)
                service.detailed_content = json.dumps(content)
                db.session.commit()
                generated += 1
            except Exception as e:
                print(f"Error generating content for service {service.name}: {e}")

    if item_type in ["all", "manufacturing"]:
        products = ManufacturingProduct.query.filter(
            (ManufacturingProduct.detailed_content == None)
            | (ManufacturingProduct.detailed_content == "")
        ).all()
        for product in products:
            try:
                prompt = f"""You are a professional content writer for Trident Engineers & Associates, an electrical panel and industrial manufacturing company.
Generate detailed content for: {product.name}
Description: {product.short_description}
Category: {product.category}

Respond ONLY in this exact JSON format:
{{"tagline": "compelling tagline", "overview": "3-4 detailed paragraphs", "key_features": ["feature1", "feature2", "feature3", "feature4", "feature5", "feature6"], "specifications": [{{"spec": "name", "value": "value"}}, {{"spec": "name", "value": "value"}}, {{"spec": "name", "value": "value"}}, {{"spec": "name", "value": "value"}}, {{"spec": "name", "value": "value"}}], "applications": ["app1", "app2", "app3", "app4", "app5"], "manufacturing_process": [{{"step": "title", "description": "desc"}}, {{"step": "title", "description": "desc"}}, {{"step": "title", "description": "desc"}}, {{"step": "title", "description": "desc"}}], "quality_standards": "paragraph about standards", "why_choose_us": ["reason1", "reason2", "reason3", "reason4"]}}"""

                response = groq_client.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role": "user", "content": prompt}],
                )
                text_resp = (
                    response.choices[0]
                    .message.content.strip()
                    .replace("```json", "")
                    .replace("```", "")
                )
                content = json.loads(text_resp)
                product.detailed_content = json.dumps(content)
                db.session.commit()
                generated += 1
            except Exception as e:
                print(
                    f"Error generating content for product {product.name}: {e}"
                )

    return redirect(request.referrer or "/services")


def _replace_text_in_slide(slide, replacements):
    def process_shapes(shapes):
        for shape in shapes:
            # Recurse into groups
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                process_shapes(shape.shapes)
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # Merge all runs to handle split placeholders
                full_text = "".join(run.text for run in para.runs)
                replaced = False
                for placeholder, value in replacements.items():
                    if placeholder in full_text:
                        full_text = full_text.replace(placeholder, value)
                        replaced = True
                if replaced and para.runs:
                    para.runs[0].text = full_text
                    for run in para.runs[1:]:
                        run.text = ""
                elif not replaced:
                    # Per-run replacement for non-split cases
                    for run in para.runs:
                        for placeholder, value in replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, value)

    process_shapes(slide.shapes)

def shorten_title(text, max_length=28):

    if not text:

        return ""

    text = text.strip()

    if len(text) <= max_length:

        return text

    words = text.split()

    result = ""

    for word in words:

        candidate = result + " " + word if result else word

        if len(candidate) > max_length - 3:

            break

        result = candidate

    return result + "..."

def _replace_slide_image(slide, image_url, target_index=1):
    import requests

    try:

        response = requests.get(image_url, timeout=10)

        response.raise_for_status()

        image_bytes = response.content

        pictures = []

        for shape in slide.shapes:

            if shape.shape_type == 13:   # Picture

                area = shape.width * shape.height

                pictures.append((area, shape))

        if not pictures:

            print("No pictures found.")

            return

        # Ignore the logo by choosing the largest picture

        pictures.sort(reverse=True)

        picture = pictures[0][1]

        blips = picture._element.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
        )

        if not blips:

            return

        r_embed_attr = (
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
        )

        rId = blips[0].get(r_embed_attr)

        if rId and rId in slide.part.rels:

            slide.part.rels[rId].target_part._blob = image_bytes

            print("Image replaced.")

    except Exception as e:

        print(e)

@app.route("/generate-service-ppt/<int:id>")
def generate_service_ppt(id):
    """Generate a PPT for a single service using the Brand Strategy template (Slide 6)."""
    service = Service.query.get_or_404(id)

    content = None
    if service.detailed_content:
        try:
            content = json.loads(service.detailed_content)
        except:
            content = None

    service_name = service.name
    overview_text = content["overview"] if content and content.get("overview") else (service.short_description or "")
    services_text = "\n".join([f"- {f}" for f in content["key_features"]]) if content and content.get("key_features") else (service.short_description or "")
    why_text = "\n".join([f"- {r}" for r in content["why_choose_us"]]) if content and content.get("why_choose_us") else "- Professional team\n- Quality workmanship\n- Indian standards compliance\n- End-to-end project management"

    template_path = os.path.join("templates", "Brand Strategy.pptx")
    prs = Presentation(template_path)
    slide = prs.slides[5]  # Slide 6 = service template

    _replace_text_in_slide(slide, {
        "{{SERVICE_NAME}}": service_name,
        "{{SERVICE_OVERVIEW}}": overview_text,
        "{{THE_SERVICES}}": services_text,
        "{{WHY_CHOOSE_US}}": why_text,
    })

    #if service.image_url:
        #replace_slide_image(slide, service.image_url)

    # Clear manufacturing placeholders on slide 7 so they don't appear raw
    _replace_text_in_slide(prs.slides[6], {
        "{{MANUFACTURING_NAME}}": "",
        "{{MANUFACTURING_OVERVIEW}}": "",
    })

    safe_name = service.name.replace("/", "-").replace("\\", "-").replace(":", "-").strip()
    filename = f"{safe_name}_Service.pptx"
    output_path = os.path.join("static", "uploads", filename)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return send_file(output_path, as_attachment=True, download_name=filename)


@app.route("/generate-manufacturing-ppt/<int:id>")
def generate_manufacturing_ppt(id):
    """Generate a PPT for a single manufacturing product using the Brand Strategy template (Slide 7)."""
    import requests as req_lib

    product = ManufacturingProduct.query.get_or_404(id)

    content = None
    if product.detailed_content:
        try:
            content = json.loads(product.detailed_content)
        except:
            content = None

    mfg_name = product.name
    overview_text = content["overview"] if content and content.get("overview") else (product.short_description or "")

    template_path = os.path.join("templates", "Brand Strategy.pptx")
    prs = Presentation(template_path)
    slide = prs.slides[6]  # Slide 7 = manufacturing template

    _replace_text_in_slide(slide, {
        "{{MANUFACTURING_NAME}}": mfg_name,
        "{{MANUFACTURING_OVERVIEW}}": overview_text,
    })

    #if product.image_url:
    #    _replace_slide_image(slide, product.image_url, target_index=1)

    # Clear service placeholders on slide 6
    _replace_text_in_slide(prs.slides[5], {
        "{{SERVICE_NAME}}": "",
        "{{SERVICE_OVERVIEW}}": "",
        "{{THE_SERVICES}}": "",
        "{{WHY_CHOOSE_US}}": "",
    })

    safe_name = product.name.replace("/", "-").replace("\\", "-").replace(":", "-").strip()
    filename = f"{safe_name}_Manufacturing.pptx"
    output_path = os.path.join("static", "uploads", filename)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return send_file(output_path, as_attachment=True, download_name=filename)

@app.route("/generate-custom-ppt", methods=["POST"])
def generate_custom_ppt():

    import aspose.slides as slides

    service_ids = request.form.getlist("service_ids")

    product_ids = request.form.getlist("product_ids")

    if not service_ids and not product_ids:

        return redirect(
            request.referrer or "/build-ppt"
        )

    template_path = os.path.join(
        "templates",
        "Brand Strategy.pptx"
    )

    pres = slides.Presentation(template_path)

    SERVICE_TEMPLATE = 5
    MANUFACTURING_TEMPLATE = 6

    insert_position = 7

    for sid in service_ids:

        service = Service.query.get(int(sid))

        if not service:

            continue

        content = None

        if service.detailed_content:

            try:

                content = json.loads(
                    service.detailed_content
                )

            except:

                content = None

        new_slide = pres.slides.insert_clone(
            insert_position,
            pres.slides[SERVICE_TEMPLATE]
        )

        insert_position += 1

        overview = (
            content["overview"]

            if content and content.get("overview")

            else service.short_description or ""
        )

        services = (
            "\n".join(
                [
                    f"• {x}"

                    for x in content["key_features"]
                ]
            )

            if content and content.get("key_features")

            else ""
        )

        why = (
            "\n".join(
                [
                    f"✓ {x}"

                    for x in content["why_choose_us"]
                ]
            )

            if content and content.get("why_choose_us")

            else ""
        )

        replacements = {

            "{{SERVICE_NAME}}":
            shorten_title(service.name),

            "{{SERVICE_OVERVIEW}}":
            overview,

            "{{THE_SERVICES}}":
            services,

            "{{WHY_CHOOSE_US}}":
            why,

        }

        for shape in new_slide.shapes:

            if hasattr(shape, "text_frame"):

                if shape.text_frame:

                    text = shape.text_frame.text

                    for old, new in replacements.items():

                        text = text.replace(
                            old,
                            new
                        )

                    shape.text_frame.text = text

    for pid in product_ids:

        product = ManufacturingProduct.query.get(
            int(pid)
        )

        if not product:

            continue

        content = None

        if product.detailed_content:

            try:

                content = json.loads(
                    product.detailed_content
                )

            except:

                content = None

        new_slide = pres.slides.insert_clone(
            insert_position,
            pres.slides[MANUFACTURING_TEMPLATE]
        )

        insert_position += 1

        overview = (
            content["overview"]

            if content and content.get("overview")

            else product.short_description or ""
        )

        replacements = {

            "{{MANUFACTURING_NAME}}":
            shorten_title(product.name),

            "{{MANUFACTURING_OVERVIEW}}":
            overview,

        }

        for shape in new_slide.shapes:

            if hasattr(shape, "text_frame"):

                if shape.text_frame:

                    text = shape.text_frame.text

                    for old, new in replacements.items():

                        text = text.replace(
                            old,
                            new
                        )

                    shape.text_frame.text = text

    pres.slides.remove_at(
        MANUFACTURING_TEMPLATE
    )

    pres.slides.remove_at(
        SERVICE_TEMPLATE
    )

    filename = (
        "Trident_Custom_Selection.pptx"
    )

    output_path = os.path.join(

        "static",

        "uploads",

        filename

    )

    pres.save(

        output_path,

        slides.export.SaveFormat.PPTX

    )

    return send_file(

        output_path,

        as_attachment=True,

        download_name=filename

    )
@app.route("/build-ppt")
def build_ppt_picker():

    all_services = Service.query.all()

    all_products = ManufacturingProduct.query.all()

    return render_template(
        "build_ppt.html",
        services=all_services,
        products=all_products
    )
@app.route("/test-ppt")
def test_ppt():

    template_path = os.path.join(
        "templates",
        "Brand Strategy.pptx"
    )

    prs = Presentation(template_path)

    output_path = os.path.join(
        "static",
        "uploads",
        "test_output.pptx"
    )

    prs.save(output_path)

    return send_file(
        output_path,
        as_attachment=True
    )
if __name__ == "__main__":
    with app.app_context():

        db.create_all()

    app.run(
        debug=True,
        use_reloader=False
    )
