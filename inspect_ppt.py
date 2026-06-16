"""Inspect Group 2 on slide 6 to understand the image structure."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

prs = Presentation("templates/Brand Strategy.pptx")
slide = prs.slides[5]  # Slide 6

ns = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
}

# Find Group 2
for s in slide.shapes:
    if s.name == 'Group 2':
        print(f"Group 2 found: type={s.shape_type}")
        print(f"  Position: ({Emu(s.left).inches:.2f}in, {Emu(s.top).inches:.2f}in)")
        print(f"  Size: ({Emu(s.width).inches:.2f}in x {Emu(s.height).inches:.2f}in)")
        
        # Get the pic element
        pics = s._element.findall('.//p:pic', ns)
        print(f"  Found {len(pics)} pic elements")
        
        for pic in pics:
            blip_fill = pic.find('.//a:blipFill', ns)
            if blip_fill is not None:
                blip = blip_fill.find('a:blip', ns)
                if blip is not None:
                    rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    print(f"  Image rId: {rId}")
                    # Get the relationship
                    rel = slide.part.rels[rId]
                    print(f"  Target: {rel.target_ref}")
                    print(f"  Image content type: {rel.target_part.content_type}")
                    print(f"  Image blob size: {len(rel.target_part.blob)} bytes")
            
            # Get size of the pic inside group
            spPr = pic.find('.//a:xfrm', ns)
            if spPr is not None:
                off = spPr.find('a:off', ns)
                ext = spPr.find('a:ext', ns)
                if off is not None and ext is not None:
                    print(f"  Pic offset in group: x={off.get('x')}, y={off.get('y')}")
                    print(f"  Pic size in group: cx={ext.get('cx')}, cy={ext.get('cy')}")
        
        # Print group XML summary
        xml = etree.tostring(s._element, pretty_print=True).decode()
        print(f"\n  Full Group XML ({len(xml)} chars):")
        print(xml[:2000])

# Also check slide 7 for manufacturing image
print("\n\n=== Slide 7 (Manufacturing) ===")
slide7 = prs.slides[6]
for s in slide7.shapes:
    if s.shape_type == 6 or 'Group' in s.name or 'Picture' in s.name:  # GROUP type
        print(f"Shape: name='{s.name}', type={s.shape_type}")
        pics = s._element.findall('.//p:pic', ns)
        print(f"  Found {len(pics)} pic elements")
        for pic in pics:
            blip = pic.find('.//a:blip', ns)
            if blip is not None:
                rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                print(f"  Image rId: {rId}")

# List all shapes on slide 7
print("\nAll slide 7 shapes:")
for s in slide7.shapes:
    print(f"  name='{s.name}', type={s.shape_type}")
    if hasattr(s._element, 'findall'):
        pics = s._element.findall('.//p:pic', ns) 
        blips = s._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill')
        if pics or blips:
            print(f"    ** Has {len(pics)} pics, {len(blips)} blipFills **")
